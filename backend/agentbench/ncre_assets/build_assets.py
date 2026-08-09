"""一次性素材重建脚本：从文字/CSV 数据重建 NCRE 真题素材并生成 blobs 分卷模块。

原卷二进制素材未抓到（详见 .ncre-sources/materials/source_notes.md 与
README_paper02_03.md），按数据稿口径用 python-docx / openpyxl / python-pptx
重建等价素材：
- paper01：Word.docx（邀请函文稿）+ Excel.xlsx（订单三表）+ 图书策划方案.docx
  → blobs.py（常量不带前缀，保持原行为）。
- paper02：WORD.docx（海报文稿）+ Excel.xlsx（成绩表）+ 水资源素材大纲.docx
  → blobs_paper02.py（常量带 PAPER02_ 前缀）。
- paper03：年报.docx（统计工作年报）+ 两份普查 CSV + 两份物理课件源 pptx
  → blobs_paper03.py（常量带 PAPER03_ 前缀）。

用法：.venv python -m agentbench.ncre_assets.build_assets [paper01|paper02|paper03]
产物：同目录 blobs.py / blobs_paper02.py / blobs_paper03.py（base64 常量，
供 catalog.py 嵌入 initial_files）。
"""
from __future__ import annotations

import base64
import datetime as dt
import io
import re
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font

from .exam_data import (
    BOOK_REF,
    CONTACTS_CSV,
    ORDER_ROWS,
    PPT_DOC_STRUCTURE,
    STATS_LABELS,
    WORD_PARAGRAPHS,
)


def build_word_asset() -> bytes:
    """邀请函素材：7 段纯文字，未做任何格式设置。"""
    document = Document()
    for text in WORD_PARAGRAPHS:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_excel_asset() -> bytes:
    """销售数据素材：三工作表，待填列留空。"""
    workbook = Workbook()
    orders = workbook.active
    orders.title = "订单明细表"
    headers = ["订单编号", "订单日期", "客户名称", "图书编号", "图书名称", "单价", "销量(本)", "小计"]
    for column, header in enumerate(headers, start=1):
        cell = orders.cell(row=2, column=column, value=header)
        cell.font = Font(bold=True)
    for offset, (order_no, date_text, customer, book_no, quantity) in enumerate(ORDER_ROWS):
        row = 3 + offset
        orders.cell(row=row, column=1, value=order_no)
        date_cell = orders.cell(
            row=row, column=2, value=dt.datetime.strptime(date_text, "%Y-%m-%d")
        )
        date_cell.number_format = "yyyy-mm-dd"
        orders.cell(row=row, column=3, value=customer)
        orders.cell(row=row, column=4, value=book_no)
        # 第 5/6/8 列（图书名称、单价、小计）留空待考生填充
        orders.cell(row=row, column=7, value=quantity)

    book_ref = workbook.create_sheet("编号对照")
    book_ref["A1"] = "图书编号"
    book_ref["B1"] = "图书名称"
    book_ref["C1"] = "单价"
    for cell in book_ref[1]:
        cell.font = Font(bold=True)
    for offset, (book_no, name, price) in enumerate(BOOK_REF):
        row = 2 + offset
        book_ref.cell(row=row, column=1, value=book_no)
        book_ref.cell(row=row, column=2, value=name)
        book_ref.cell(row=row, column=3, value=price)

    report = workbook.create_sheet("统计报告")
    report["A2"] = "统计项目"
    report["B2"] = "统计结果"
    report["A2"].font = Font(bold=True)
    report["B2"].font = Font(bold=True)
    for offset, label in enumerate(STATS_LABELS):
        report.cell(row=3 + offset, column=1, value=label)

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_ppt_source_asset() -> bytes:
    """图书策划方案素材：按标题 1/2/3 层级重建（正文层为 Heading 3）。"""
    document = Document()
    for style, text in PPT_DOC_STRUCTURE:
        document.add_paragraph(text, style=style)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------- paper02
def build_paper02_word_asset() -> bytes:
    """海报素材：第1页海报文字 + 第2页标题与日程文字，未做任何格式设置。"""
    from .exam_data_paper02 import (
        SCHEDULE_HEADER,
        SCHEDULE_ROWS,
        SPEAKER_INTRO,
        WORD_PAGE2_TITLE,
        WORD_POSTER_PAGE1,
    )

    document = Document()
    for text in WORD_POSTER_PAGE1:
        document.add_paragraph(text)
    document.add_paragraph(WORD_PAGE2_TITLE)
    document.add_paragraph("日程安排：" + "、".join(SCHEDULE_HEADER))
    for time_range, topic, speaker in SCHEDULE_ROWS:
        document.add_paragraph(f"{time_range}，{topic}，{speaker}")
    document.add_paragraph("报告人介绍")
    document.add_paragraph(SPEAKER_INTRO)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_paper02_excel_asset() -> bytes:
    """成绩单素材：成绩表工作表，学号文本列 + 7科成绩，班级列留空。"""
    from .exam_data_paper02 import GRADE_ROWS, SUBJECTS

    workbook = Workbook()
    grades = workbook.active
    grades.title = "成绩表"
    headers = ["学号", "姓名", "班级"] + SUBJECTS
    for column, header in enumerate(headers, start=1):
        grades.cell(row=1, column=column, value=header).font = Font(bold=True)
    for offset, (student_id, name, _class_no, *scores) in enumerate(GRADE_ROWS):
        row = 2 + offset
        id_cell = grades.cell(row=row, column=1, value=str(student_id))
        id_cell.number_format = "@"
        grades.cell(row=row, column=2, value=name)
        # 第 3 列班级留空待考生用公式提取
        for index, score in enumerate(scores):
            grades.cell(row=row, column=4 + index, value=score)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_paper02_ppt_source_asset() -> bytes:
    """水资源利用与节水素材大纲：一/二/三级标题 + 正文。"""
    from .exam_data_paper02 import PPT_OUTLINE

    document = Document()
    for style, text in PPT_OUTLINE:
        document.add_paragraph(text, style=style)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------- paper03
def build_paper03_report_asset() -> bytes:
    """年报素材：封面/引言/正文纯文字重建；红色超链接候选文字与蓝色咨询段保留提示色。"""
    from docx.shared import RGBColor

    from .exam_data_paper03 import LINK_TEXT, REPORT_BODY, REPORT_COVER, REPORT_INTRO

    document = Document()
    for text in REPORT_COVER:
        document.add_paragraph(text)
    document.add_paragraph(REPORT_INTRO[0])
    before, after = REPORT_INTRO[1].split(LINK_TEXT)
    intro_paragraph = document.add_paragraph(before)
    link_run = intro_paragraph.add_run(LINK_TEXT)
    link_run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    intro_paragraph.add_run(after)
    for kind, text in REPORT_BODY:
        paragraph = document.add_paragraph(text)
        if kind == "blue":
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _rename_theme(payload: bytes, theme_name: str) -> bytes:
    """zip 后处理：重命名 ppt/theme/theme1.xml 的主题名。"""
    with zipfile.ZipFile(io.BytesIO(payload)) as archive:
        names = archive.namelist()
        contents = {name: archive.read(name) for name in names}
    theme_xml = contents["ppt/theme/theme1.xml"].decode("utf-8")
    theme_xml = re.sub(r'(<a:theme[^>]*name=")[^"]*(")', rf"\g<1>{theme_name}\g<2>", theme_xml)
    contents["ppt/theme/theme1.xml"] = theme_xml.encode("utf-8")
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for name in names:
            archive.writestr(name, contents[name])
    return buffer.getvalue()


def build_paper03_courseware_asset(slides_spec: list, theme_name: str) -> bytes:
    """物理课件源 pptx：按 (版式, 标题, 内容要点) 结构重建并应用指定主题名。"""
    from pptx import Presentation

    presentation = Presentation()
    for kind, title, bullets in slides_spec:
        layout = presentation.slide_layouts[0 if kind == "title" else 1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text_frame.text = title
        body = next(p for p in slide.placeholders if p.placeholder_format.idx != 0)
        frame = body.text_frame
        frame.clear()
        for offset, line in enumerate(bullets):
            paragraph = frame.paragraphs[0] if offset == 0 else frame.add_paragraph()
            paragraph.text = line
    buffer = io.BytesIO()
    presentation.save(buffer)
    return _rename_theme(buffer.getvalue(), theme_name)


def build_paper03_courseware_assets() -> dict[str, tuple[str, bytes]]:
    from .exam_data_paper03 import COURSEWARE_A_SLIDES, COURSEWARE_B_SLIDES

    return {
        "PAPER03_COURSEWARE_A_PPTX_B64": (
            "第1-2节.pptx",
            build_paper03_courseware_asset(COURSEWARE_A_SLIDES, "主题A 科技蓝"),
        ),
        "PAPER03_COURSEWARE_B_PPTX_B64": (
            "第3-5节.pptx",
            build_paper03_courseware_asset(COURSEWARE_B_SLIDES, "主题B 简约灰"),
        ),
    }


# ---------------------------------------------------------------- 输出
def _emit_blobs_module(target: Path, text_constants: dict[str, str],
                       assets: dict[str, tuple[str, bytes]]) -> None:
    lines = [
        '"""由 build_assets.py 生成的 base64 素材常量，请勿手改。"""',
        "from __future__ import annotations",
        "",
    ]
    for constant, content in text_constants.items():
        lines.append(f'{constant} = """{content}"""')
        lines.append("")
    for constant, (filename, payload) in assets.items():
        encoded = base64.b64encode(payload).decode("ascii")
        chunks = "\n    ".join(f'"{encoded[i:i + 76]}"' for i in range(0, len(encoded), 76))
        lines.append(f"# {filename}（{len(payload)} 字节）")
        lines.append(f"{constant} = (\n    {chunks}\n)")
        lines.append("")
    target.write_text("\n".join(lines), encoding="utf-8")
    print(f"{target.name} written: {target} ({target.stat().st_size} bytes)")


def main(paper: str = "paper01") -> None:
    if paper == "paper01":
        assets = {
            "WORD_DOCX_B64": ("Word.docx", build_word_asset()),
            "EXCEL_XLSX_B64": ("Excel.xlsx", build_excel_asset()),
            "PPT_SOURCE_DOCX_B64": ("图书策划方案.docx", build_ppt_source_asset()),
        }
        _emit_blobs_module(
            Path(__file__).with_name("blobs.py"),
            {"CONTACTS_CSV": CONTACTS_CSV},
            assets,
        )
    elif paper == "paper02":
        from .exam_data_paper02 import GRADES_CSV

        assets = {
            "PAPER02_WORD_DOCX_B64": ("WORD.docx", build_paper02_word_asset()),
            "PAPER02_EXCEL_XLSX_B64": ("Excel.xlsx", build_paper02_excel_asset()),
            "PAPER02_PPT_SOURCE_DOCX_B64": (
                "水资源素材大纲.docx", build_paper02_ppt_source_asset()),
        }
        _emit_blobs_module(
            Path(__file__).with_name("blobs_paper02.py"),
            {"PAPER02_GRADES_CSV": GRADES_CSV},
            assets,
        )
    elif paper == "paper03":
        from .exam_data_paper03 import (
            CENSUS_2000_CSV,
            CENSUS_2010_CSV,
            CONSULT_CSV,
        )

        assets = {
            "PAPER03_REPORT_DOCX_B64": ("年报.docx", build_paper03_report_asset()),
            **build_paper03_courseware_assets(),
        }
        _emit_blobs_module(
            Path(__file__).with_name("blobs_paper03.py"),
            {
                "PAPER03_CENSUS_2000_CSV": CENSUS_2000_CSV,
                "PAPER03_CENSUS_2010_CSV": CENSUS_2010_CSV,
                "PAPER03_CONSULT_CSV": CONSULT_CSV,
            },
            assets,
        )
    else:
        raise SystemExit(f"未知卷别: {paper}（可选 paper01/paper02/paper03）")


if __name__ == "__main__":
    import sys

    main(sys.argv[1] if len(sys.argv) > 1 else "paper01")
