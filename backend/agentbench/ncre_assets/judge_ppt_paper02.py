# NCRE PowerPoint 操作题私有判分脚本（经典题库第2套：水资源利用与节水，command_metrics 协议）
# 用 python-pptx + zipfile 检查 水资源利用与节水.pptx 的数量、主题、标题页、版式、
# 图片与超链接；产物缺失时全 0。
import contextlib
import json
import re
import sys
import zipfile

with contextlib.suppress(Exception):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

metrics = {f"p{i}": 0.0 for i in range(1, 8)}
evidence = {}

TITLE = "水资源利用与节水"
ORG = "北京节水展馆"
KEYWORDS = ["水资源", "海水淡化", "节水器具"]
FILE_NAME = "水资源利用与节水.pptx"


def emit():
    print("AGENTBENCH_METRICS=" + json.dumps({"metrics": metrics, "evidence": evidence}))


try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"python-pptx 不可用: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

presentation = None
try:
    presentation = Presentation(FILE_NAME)
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"{FILE_NAME} 缺失或无法打开: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

slides = list(presentation.slides)

# ------------------------------------------------ p1 幻灯片数量 >= 5
metrics["p1"] = 100.0 if len(slides) >= 5 else 0.0
evidence["p1"] = f"{len(slides)} 张"

# ------------------------------------------------ p2 非默认主题
archive = None
try:
    archive = zipfile.ZipFile(FILE_NAME)
    theme_xml = archive.read("ppt/theme/theme1.xml").decode("utf-8", errors="replace")
    theme_name = re.search(r'<a:theme[^>]*name="([^"]*)"', theme_xml)
    font_scheme = re.search(r'<a:fontScheme[^>]*name="([^"]*)"', theme_xml)
    color_scheme = re.search(r'<a:clrScheme[^>]*name="([^"]*)"', theme_xml)
    non_default = (theme_name and theme_name.group(1) != "Office Theme") or (
        font_scheme and font_scheme.group(1) != "Office"
    ) or (color_scheme and color_scheme.group(1) not in ("Office", "Office 主题"))
    metrics["p2"] = 100.0 if non_default else 0.0
    evidence["p2"] = theme_name.group(1) if theme_name else "?"
except Exception as exc:  # noqa: BLE001
    evidence["p2"] = str(exc)[:200]

# ------------------------------------------------ p3 标题页要素
try:
    first = slides[0]
    texts = " ".join(
        shape.text_frame.text for shape in first.shapes if shape.has_text_frame)
    layout_ok = first.slide_layout.name in ("Title Slide", "标题幻灯片")
    title_ok = TITLE in texts
    org_ok = ORG in texts
    date_ok = bool(re.search(
        r"\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日|\d{4}[-/.]\d{1,2}[-/.]\d{1,2}", texts))
    metrics["p3"] = 25.0 * sum([layout_ok, title_ok, org_ok, date_ok])
    evidence["p3"] = f"版式ok={layout_ok}, 标题ok={title_ok}, 单位ok={org_ok}, 日期ok={date_ok}"
except Exception as exc:  # noqa: BLE001
    evidence["p3"] = str(exc)[:200]

# ------------------------------------------------ p4 内容板块关键词与版式多样性
try:
    all_text = " ".join(
        shape.text_frame.text
        for slide in slides for shape in slide.shapes if shape.has_text_frame)
    keyword_hits = sum(1 for keyword in KEYWORDS if keyword in all_text)
    layouts = {slide.slide_layout.name for slide in slides}
    layout_ok = len(layouts) >= 3
    metrics["p4"] = round(60.0 * keyword_hits / len(KEYWORDS) + 40.0 * layout_ok, 2)
    evidence["p4"] = f"关键词{keyword_hits}/3, 版式{sorted(layouts)}"
except Exception as exc:  # noqa: BLE001
    evidence["p4"] = str(exc)[:200]

# ------------------------------------------------ p5 图片数量 >= 2
try:
    def count_pictures(shapes):
        total = 0
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total += 1
            if getattr(shape, "shapes", None):
                total += count_pictures(shape.shapes)
        return total

    pictures = sum(count_pictures(slide.shapes) for slide in slides)
    metrics["p5"] = 100.0 if pictures >= 2 else 100.0 * pictures / 2
    evidence["p5"] = f"{pictures} 张图片"
except Exception as exc:  # noqa: BLE001
    evidence["p5"] = str(exc)[:200]

# ------------------------------------------------ p6 超链接数量 >= 2
try:
    link_count = 0
    for name in archive.namelist():
        if not re.search(r"ppt/slides/_rels/slide\d+\.xml\.rels$", name):
            continue
        rels_xml = archive.read(name).decode("utf-8", errors="replace")
        for rel_type, _target, mode in re.findall(
                r'Type="([^"]*)"(?=(?:(?!/>).)*?Target="([^"]*)")((?:(?!/>).)*)(?=/?>)', rels_xml):
            if rel_type.endswith("/hyperlink") or rel_type.endswith("/slide") and "TargetMode" not in mode:
                link_count += 1
    metrics["p6"] = 100.0 if link_count >= 2 else 100.0 * link_count / 2
    evidence["p6"] = f"{link_count} 个超链接"
except Exception as exc:  # noqa: BLE001
    evidence["p6"] = str(exc)[:200]
finally:
    if archive is not None:
        archive.close()

# ------------------------------------------------ p7 文件名
metrics["p7"] = 100.0
evidence["p7"] = FILE_NAME

emit()
