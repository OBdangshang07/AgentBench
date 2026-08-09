# NCRE PowerPoint 操作题私有判分脚本（经典题库第4套：物理课件整合，command_metrics 协议）
# 用 python-pptx + zipfile 检查 物理课件.pptx 的双主题合并、新增页、对比表、
# 内部超链接、编号页脚与切换；产物缺失时全 0。
import contextlib
import json
import re
import sys
import zipfile

with contextlib.suppress(Exception):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

metrics = {f"p{i}": 0.0 for i in range(1, 8)}
evidence = {}

FILE_NAME = "物理课件.pptx"
TITLE_ONLY_NAMES = ("仅标题", "Title Only")
NEW_TITLES = {3: "物质的状态", 6: "蒸发和沸腾的异同点"}
TABLE_HEADER = ["项目", "蒸发", "沸腾"]
TABLE_KEYWORDS = ["液体表面", "内部和表面", "任何温度", "沸点"]
LINKS = [(3, 2, "返回第2节熔化和凝固"), (6, 5, "返回第4节升华和凝华")]
FOOTER_TEXT = "第一章 物态及其变化"


def emit():
    print("AGENTBENCH_METRICS=" + json.dumps({"metrics": metrics, "evidence": evidence}))


try:
    from pptx import Presentation
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"python-pptx 不可用: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

presentation = None
try:
    presentation = Presentation(FILE_NAME)
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"{FILE_NAME} 缺失或无法打开: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

slides = list(presentation.slides)

# ------------------------------------------------ p1 9页合并与双主题
try:
    count_ok = len(slides) == 9
    masters = {id(slide.slide_layout.slide_master) for slide in slides}
    metrics["p1"] = round(50.0 * count_ok + 50.0 * (len(masters) >= 2), 2)
    evidence["p1"] = f"{len(slides)} 张, {len(masters)} 个主题"
except Exception as exc:  # noqa: BLE001
    evidence["p1"] = str(exc)[:200]


def slide_title(slide):
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


# ------------------------------------------------ p2 物质的状态（仅标题版式）
try:
    if len(slides) > 3:
        layout_name = slides[3].slide_layout.name
        layout_ok = any(name in layout_name for name in TITLE_ONLY_NAMES)
        title_ok = slide_title(slides[3]) == NEW_TITLES[3]
        metrics["p2"] = 50.0 * layout_ok + 50.0 * title_ok
        evidence["p2"] = f"版式={layout_name}, 标题ok={title_ok}"
    else:
        evidence["p2"] = "幻灯片不足4张"
except Exception as exc:  # noqa: BLE001
    evidence["p2"] = str(exc)[:200]

# ------------------------------------------------ p3 蒸发沸腾异同表格页
try:
    if len(slides) > 6:
        for shape in slides[6].shapes:
            if not shape.has_table:
                continue
            table = shape.table
            header = [cell.text.strip() for cell in table.rows[0].cells]
            structure_ok = (len(table.rows) == 4 and len(table.columns) == 3
                            and header == TABLE_HEADER
                            and slide_title(slides[6]) == NEW_TITLES[6])
            content = " ".join(
                cell.text for row in table.rows for cell in row.cells)
            keyword_hits = sum(1 for keyword in TABLE_KEYWORDS if keyword in content)
            metrics["p3"] = round(50.0 * structure_ok
                                  + 50.0 * keyword_hits / len(TABLE_KEYWORDS), 2)
            evidence["p3"] = f"结构ok={structure_ok}, 关键词{keyword_hits}/4"
            break
        else:
            evidence["p3"] = "第7张无表格"
    else:
        evidence["p3"] = "幻灯片不足7张"
except Exception as exc:  # noqa: BLE001
    evidence["p3"] = str(exc)[:200]

# ------------------------------------------------ p4-p6 基于包内 XML
archive = None
try:
    archive = zipfile.ZipFile(FILE_NAME)
    slide_parts = [str(slide.part.partname).lstrip("/") for slide in slides]

    def slide_xml(index):
        return archive.read(slide_parts[index]).decode("utf-8", errors="replace")

    def slide_rels(index):
        part = slide_parts[index]
        rels_name = part.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
        if rels_name not in archive.namelist():
            return ""
        return archive.read(rels_name).decode("utf-8", errors="replace")

    # p4 内部超链接：第4张→第3张、第7张→第6张
    score = 0.0
    detail = []
    for source, target_index, link_text in LINKS:
        if source >= len(slides):
            detail.append(f"缺第{source + 1}张")
            continue
        expected_target = slide_parts[target_index].rsplit("/", 1)[-1]
        rels_xml = slide_rels(source)
        link_ok = False
        for rel_type, target in re.findall(r'Type="([^"]*)"[^>]*Target="([^"]*)"', rels_xml):
            if rel_type.endswith("/slide") and target.rsplit("/", 1)[-1] == expected_target:
                link_ok = True
                break
        xml = slide_xml(source)
        text_ok = link_text in xml
        score += 25.0 * link_ok + 25.0 * text_ok
        detail.append(f"第{source + 1}张: 链接={link_ok}, 文字={text_ok}")
    metrics["p4"] = round(score, 2)
    evidence["p4"] = "; ".join(detail)

    # p5 页脚与幻灯片编号
    footer_hits = number_hits = 0
    for index in range(len(slides)):
        xml = slide_xml(index)
        if '<p:ph type="ftr"' in xml and FOOTER_TEXT in xml:
            footer_hits += 1
        if index > 0 and '<p:ph type="sldNum"' in xml:
            number_hits += 1
    metrics["p5"] = round(
        60.0 * footer_hits / max(len(slides), 1)
        + 40.0 * number_hits / max(len(slides) - 1, 1), 2)
    evidence["p5"] = f"页脚{footer_hits}/{len(slides)}, 编号{number_hits}/{max(len(slides) - 1, 0)}"

    # p6 幻灯片切换
    transitions = []
    for index in range(len(slides)):
        match = re.search(r"<p:transition[^>]*>\s*(<p:(\w+))?", slide_xml(index))
        if match and "<p:transition" in slide_xml(index):
            transitions.append(match.group(2) or "none")
    present = len(transitions)
    consistent = len(set(transitions)) == 1 if transitions else False
    metrics["p6"] = round(50.0 * present / max(len(slides), 1) + 50.0 * consistent, 2)
    evidence["p6"] = f"{present}/{len(slides)} 页有切换, 一致={consistent}"
except Exception as exc:  # noqa: BLE001
    evidence.setdefault("p4", str(exc)[:200])
finally:
    if archive is not None:
        archive.close()

# ------------------------------------------------ p7 文件名
metrics["p7"] = 100.0
evidence["p7"] = FILE_NAME

emit()
