# NCRE PowerPoint 操作题私有判分脚本（AgentBench command_metrics 协议）
# 用 python-pptx + zipfile 检查 PowerPoint.pptx 的结构、主题、表格、SmartArt 与自定义放映。
import contextlib
import json
import re
import sys
import zipfile

with contextlib.suppress(Exception):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

metrics = {f"p{i}": 0.0 for i in range(1, 9)}
evidence = {}

TITLES = [
    "Microsoft Office图书策划案",
    "推荐作者简介",
    "Office 2010的十大优势",
    "新版图书读者定位",
    "PowerPoint 2010创新的功能体验",
    "2012年同类图书销量统计",
    "新版图书创作流程示意",
]
FLOW_STEPS = ["确定选题", "选定作者", "选题沟通", "图书编写", "编辑审校", "排版印刷", "上市发行"]
TABLE_HEADERS = ["图书名称", "出版社", "作者", "定价", "销量"]


def emit():
    print("AGENTBENCH_METRICS=" + json.dumps({"metrics": metrics, "evidence": evidence}))


try:
    from pptx import Presentation
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"python-pptx 不可用: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

presentation = None
try:
    presentation = Presentation("PowerPoint.pptx")
except Exception as exc:  # noqa: BLE001
    evidence["error"] = f"PowerPoint.pptx 缺失或无法打开: {str(exc)[:200]}"
    emit()
    raise SystemExit(0) from None

slides = list(presentation.slides)

# ------------------------------------------------ p1 幻灯片数量
metrics["p1"] = 100.0 if len(slides) == 7 else 0.0
evidence["p1"] = f"{len(slides)} 张"

# ------------------------------------------------ p2 各页标题
def slide_title(slide):
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
            return shape.text_frame.text.strip()
    return ""


matched = sum(
    1 for index, expected in enumerate(TITLES)
    if index < len(slides) and slide_title(slides[index]) == expected
)
metrics["p2"] = round(100.0 * matched / len(TITLES), 2)

# ------------------------------------------------ p3 层级内容抽查（第2页与第4页）
def body_paragraphs(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0 and shape.has_text_frame:
            return [(p.text.strip(), p.level) for p in shape.text_frame.paragraphs if p.text.strip()]
    for shape in slide.shapes:
        if shape.has_text_frame and shape is not slide.shapes.title:
            return [(p.text.strip(), p.level) for p in shape.text_frame.paragraphs if p.text.strip()]
    return []


score = 0.0
if len(slides) >= 4:
    page2 = body_paragraphs(slides[1])
    level0_page2 = {text for text, level in page2 if level == 0}
    if "刘雅汶" in level0_page2 and "主要代表作品" in level0_page2:
        score += 40.0
    if any(level == 1 and "Contoso" in text for text, level in page2) \
            and any(level == 1 and "Microsoft Office整合应用精要" in text for text, level in page2):
        score += 20.0
    page4 = body_paragraphs(slides[3])
    level0_page4 = {text for text, level in page4 if level == 0}
    expected_page4 = {"信息工作者", "学生和教师", "办公应用技能培训班", "大专院校教材"}
    score += 40.0 * len(expected_page4 & level0_page4) / len(expected_page4)
metrics["p3"] = round(score, 2)

# ------------------------------------------------ p4 第1页版式为标题幻灯片
try:
    layout_name = slides[0].slide_layout.name
    metrics["p4"] = 100.0 if layout_name in ("Title Slide", "标题幻灯片") else 0.0
    evidence["p4"] = layout_name
except Exception as exc:  # noqa: BLE001
    evidence["p4"] = str(exc)[:200]

# ------------------------------------------------ p5-p8 基于包内 XML
archive = None
try:
    archive = zipfile.ZipFile("PowerPoint.pptx")
    theme_xml = archive.read("ppt/theme/theme1.xml").decode("utf-8", errors="replace")
    theme_name = re.search(r'<a:theme[^>]*name="([^"]*)"', theme_xml)
    font_scheme = re.search(r'<a:fontScheme[^>]*name="([^"]*)"', theme_xml)
    color_scheme = re.search(r'<a:clrScheme[^>]*name="([^"]*)"', theme_xml)
    non_default = (theme_name and theme_name.group(1) != "Office Theme") or (
        font_scheme and font_scheme.group(1) != "Office"
    ) or (color_scheme and color_scheme.group(1) not in ("Office", "Office 主题"))
    metrics["p5"] = 100.0 if non_default else 0.0
    evidence["p5"] = theme_name.group(1) if theme_name else "?"
except Exception as exc:  # noqa: BLE001
    evidence["p5"] = str(exc)[:200]

# p6 第6页 6x5 表格
try:
    if len(slides) >= 6:
        for shape in slides[5].shapes:
            if shape.has_table:
                table = shape.table
                header_texts = [cell.text.strip() for cell in table.rows[0].cells]
                if len(table.rows) == 6 and len(table.columns) == 5 \
                        and header_texts == TABLE_HEADERS:
                    metrics["p6"] = 100.0
                evidence["p6"] = f"{len(table.rows)}x{len(table.columns)} {header_texts}"
                break
except Exception as exc:  # noqa: BLE001
    evidence["p6"] = str(exc)[:200]

# p7 第7页 SmartArt（diagram 部件 + 流程关键词）
try:
    slide7_rels = next(
        (n for n in archive.namelist() if re.search(r"ppt/slides/_rels/slide7\.xml\.rels$", n)),
        None,
    )
    if slide7_rels:
        rels_xml = archive.read(slide7_rels).decode("utf-8", errors="replace")
        diagram_parts = [
            rel_type.rsplit("/", 1)[-1]
            for rel_type in re.findall(r'Type="([^"]*)"', rels_xml)
            if re.search(r"diagram(Data|Layout|QuickStyle|Colors|Style|Drawing)$", rel_type)
        ]
        slide7_xml = archive.read("ppt/slides/slide7.xml").decode("utf-8", errors="replace")
        has_diagram_ref = bool(re.search(r"dgm:", slide7_xml))
        data_parts = [n for n in archive.namelist() if re.search(r"diagrams/data\d+\.xml$", n)]
        step_hits = 0
        if data_parts:
            data_xml = "".join(
                archive.read(part).decode("utf-8", errors="replace") for part in data_parts
            )
            step_hits = sum(1 for step in FLOW_STEPS if step in data_xml)
        metrics["p7"] = round(
            50.0 * min(len(set(diagram_parts)), 4) / 4 + 50.0 * step_hits / len(FLOW_STEPS), 2
        )
        evidence["p7"] = (
            f"diagram部件{sorted(set(diagram_parts))}, dgm引用={has_diagram_ref}, 关键词{step_hits}/7"
        )
except Exception as exc:  # noqa: BLE001
    evidence["p7"] = str(exc)[:200]

# p8 自定义放映方案
try:
    presentation_xml = archive.read("ppt/presentation.xml").decode("utf-8", errors="replace")
    rels_xml = archive.read("ppt/_rels/presentation.xml.rels").decode("utf-8", errors="replace")
    rel_to_slide = {}
    for rel_id, target in re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"', rels_xml):
        slide_match = re.search(r"slides/slide(\d+)\.xml", target)
        if slide_match:
            rel_to_slide[rel_id] = int(slide_match.group(1))
    shows = {}
    for block in re.findall(r"<p14:custShow\b[^>]*>.*?</p14:custShow>", presentation_xml, re.DOTALL):
        name = re.search(r'name="([^"]*)"', block)
        if not name:
            continue
        refs = re.findall(r'r:id="([^"]+)"', block)
        shows[name.group(1)] = sorted(rel_to_slide[r] for r in refs if r in rel_to_slide)
    expected_shows = {"放映方案1": [1, 2, 4, 7], "放映方案2": [1, 2, 3, 5, 6]}
    hits = sum(1 for name, pages in expected_shows.items() if shows.get(name) == pages)
    metrics["p8"] = round(100.0 * hits / len(expected_shows), 2)
    evidence["p8"] = str(shows)[:200]
except Exception as exc:  # noqa: BLE001
    evidence["p8"] = str(exc)[:200]
finally:
    if archive is not None:
        archive.close()

emit()
