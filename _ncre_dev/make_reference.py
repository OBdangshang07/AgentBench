"""临时脚本：生成 NCRE 三卷四道题的满分标准产物，用于本地验证判分脚本。

用法：python make_reference.py [paper01|paper02|paper03]（默认 paper01）。
产物分别写入 reference-workspace / reference-workspace-paper02 /
reference-workspace-paper03。
"""
from __future__ import annotations

import io
import json
import re
import shutil
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "backend"))

DEV = Path(__file__).parent

import openpyxl
from agentbench.ncre_assets import build_assets
from docx import Document
from docx.enum.section import WD_ORIENT, WD_SECTION_START
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from openpyxl.formatting.rule import CellIsRule
from openpyxl.styles import Font, PatternFill
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.util import Inches


def _zip_patch(path: Path, mutate) -> None:
    """zip 后处理：读全部部件 → mutate(contents) → 重写（新增键追加）。"""
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        contents = {name: archive.read(name) for name in names}
    mutate(contents)
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, payload in contents.items():
            archive.writestr(name, payload)
    path.write_bytes(buffer.getvalue())


def _add_background(docx_path: Path, color: str) -> None:
    """在 document.xml 顶层插入 w:background 纯色背景。"""

    def mutate(contents):
        xml = contents["word/document.xml"].decode("utf-8")
        xml = xml.replace("<w:body>", f'<w:background w:color="{color}"/><w:body>', 1)
        contents["word/document.xml"] = xml.encode("utf-8")

    _zip_patch(docx_path, mutate)


# ---------------------------------------------------------------- paper01
def make_paper01(out: Path) -> None:
    from agentbench.ncre_assets.exam_data import CONTACTS, WORD_PARAGRAPHS

    answers = {
        "q01": "B", "q02": "D", "q03": "C", "q04": "A", "q05": "B",
        "q06": "A", "q07": "A", "q08": "A", "q09": "A", "q10": "A",
        "q11": "D", "q12": "C", "q13": "C", "q14": "D", "q15": "B",
        "q16": "A", "q17": "B", "q18": "A", "q19": "C", "q20": "D",
    }
    (out / "answers.json").write_text(json.dumps(answers), encoding="utf-8")

    def format_template(doc: Document) -> None:
        section = doc.sections[0]
        section.page_width = Cm(30)
        section.page_height = Cm(18)
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(3)
        section.right_margin = Cm(3)
        paragraphs = doc.paragraphs
        for paragraph in paragraphs:
            for run in paragraph.runs:
                run.font.name = "微软雅黑"
                run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        for index in (0, 1):
            for run in paragraphs[index].runs:
                run.font.size = Pt(26)
            fmt = paragraphs[index].paragraph_format
            fmt.alignment = WD_ALIGN_PARAGRAPH.CENTER
            fmt.space_before = Pt(13)
            fmt.space_after = Pt(13)
        for run in paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0, 0, 255)
        for index in (3, 4):
            paragraphs[index].paragraph_format.first_line_indent = Pt(21)
        for index in (5, 6):
            paragraphs[index].alignment = WD_ALIGN_PARAGRAPH.RIGHT

    template = Document()
    for text in WORD_PARAGRAPHS:
        template.add_paragraph(text)
    format_template(template)
    template_path = out / "Word.docx"
    template.save(template_path)
    _add_background(template_path, "FDE9D9")

    merged = Document()
    merged_section = merged.sections[0]
    merged_section.page_width = Cm(30)
    merged_section.page_height = Cm(18)
    for offset, (name, _title) in enumerate(CONTACTS):
        if offset > 0:
            merged.paragraphs[-1].runs[-1].add_break(WD_BREAK.PAGE)
        for index, text in enumerate(WORD_PARAGRAPHS):
            if index == 2:
                text = f"尊敬的{name}："
            merged.add_paragraph(text)
    merged.save(out / "Word-邀请函.docx")

    workbook = openpyxl.load_workbook(io.BytesIO(build_assets.build_excel_asset()))
    orders = workbook["订单明细表"]
    last_row = 2 + 23
    orders.add_table(Table(displayName="Table1", ref=f"A2:H{last_row}"))
    orders.tables["Table1"].tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2", showRowStripes=True)
    accounting = '_ ￥* #,##0.00_ ;_ ￥* -#,##0.00_ ;_ ￥* "-"??_ ;_ @_'
    for row in range(3, last_row + 1):
        orders.cell(row=row, column=5).value = f"=VLOOKUP(D{row},编号对照!$A$2:$C$9,2,0)"
        price_cell = orders.cell(row=row, column=6)
        price_cell.value = f"=VLOOKUP(D{row},编号对照!$A$2:$C$9,3,0)"
        price_cell.number_format = accounting
        subtotal = orders.cell(row=row, column=8)
        subtotal.value = f"=F{row}*G{row}"
        subtotal.number_format = accounting
    report = workbook["统计报告"]
    report["B3"] = "=SUM(订单明细表!H3:H25)"
    report["B4"] = ('=SUMPRODUCT((订单明细表!$D$3:$D$25="BK-83021")*'
                    '(订单明细表!$B$3:$B$25>="2012-01-01")*'
                    '(订单明细表!$B$3:$B$25<="2012-12-31")*(订单明细表!$H$3:$H$25))')
    report["B5"] = ('=SUMPRODUCT((订单明细表!$C$3:$C$25="隆华书店")*'
                    '(订单明细表!$B$3:$B$25>="2011-07-01")*'
                    '(订单明细表!$B$3:$B$25<="2011-09-30")*(订单明细表!$H$3:$H$25))')
    report["B6"] = ('=SUMPRODUCT((订单明细表!$C$3:$C$25="隆华书店")*'
                    '(订单明细表!$B$3:$B$25>="2011-01-01")*'
                    '(订单明细表!$B$3:$B$25<="2011-12-31")*(订单明细表!$H$3:$H$25))/12')
    report["B6"].number_format = "0.00"
    workbook.save(out / "Excel.xlsx")

    from agentbench.ncre_assets.exam_data import (
        PPT_FLOW_STEPS,
        PPT_SLIDE_TITLES,
    )

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    prs.slides.add_slide(title_layout)
    for _ in range(6):
        prs.slides.add_slide(content_layout)
    slides = list(prs.slides)
    for slide, title in zip(slides, PPT_SLIDE_TITLES):
        slide.shapes.title.text_frame.text = title

    page2_body = [
        (0, "刘雅汶"),
        (1, "Contoso公司技术经理，微软特邀资深顾问讲师，微软全球最有价值专家（MVP）。"),
        (0, "主要代表作品"),
        (1, "《Microsoft Office整合应用精要》"),
        (1, "《Microsoft Word企业应用宝典》"),
    ]
    page4_body = [
        (0, "信息工作者"),
        (1, "应用Office产品处理日常事务。"),
        (0, "学生和教师"),
        (0, "办公应用技能培训班"),
        (0, "大专院校教材"),
    ]
    for slide, body in ((slides[1], page2_body), (slides[3], page4_body)):
        frame = next(
            p.text_frame for p in slide.placeholders if p.placeholder_format.idx != 0)
        frame.clear()
        for offset, (level, text) in enumerate(body):
            paragraph = frame.paragraphs[0] if offset == 0 else frame.add_paragraph()
            paragraph.level = level
            paragraph.text = text

    table_shape = slides[5].shapes.add_table(
        6, 5, Inches(0.8), Inches(2.0), Inches(8.4), Inches(3.0))
    table = table_shape.table
    for column, header in enumerate(["图书名称", "出版社", "作者", "定价", "销量"]):
        table.cell(0, column).text = header

    pptx_path = out / "PowerPoint.pptx"
    prs.save(pptx_path)

    def patch_pptx(path: Path) -> None:
        """zip 后处理：改主题名、注入 SmartArt 部件与自定义放映方案。"""

        def mutate(contents):
            theme = contents["ppt/theme/theme1.xml"].decode("utf-8")
            theme = re.sub(r'(<a:theme[^>]*name=")Office Theme(")', r"\1均衡\2", theme)
            contents["ppt/theme/theme1.xml"] = theme.encode("utf-8")

            steps = "".join(
                f'<dgm:pt modelId="{i + 1}"><dgm:prSet/><dgm:t>'
                f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:r><a:t>{step}</a:t></a:r></a:p></dgm:t></dgm:pt>'
                for i, step in enumerate(PPT_FLOW_STEPS)
            )
            contents["ppt/diagrams/data1.xml"] = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
                f'<dgm:ptLst><dgm:pt modelId="0" type="doc"/>{steps}</dgm:ptLst>'
                '<dgm:cxnLst/><dgm:bg/><dgm:whole/></dgm:dataModel>').encode()
            for part, root in (
                ("layout1.xml", "dgm:layoutDef"),
                ("quickStyle1.xml", "dgm:styleDef"),
                ("colors1.xml", "dgm:colorsDef"),
            ):
                contents[f"ppt/diagrams/{part}"] = (
                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                    f'<{root} xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" '
                    f'uniqueId="urn:agentbench:ncre:{part}"/>').encode()

            rels = contents["ppt/slides/_rels/slide7.xml.rels"].decode("utf-8")
            extra_rels = "".join(
                f'<Relationship Id="rIdDg{i}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/{rel_type}" '
                f'Target="../diagrams/{target}"/>'
                for i, (rel_type, target) in enumerate((
                    ("diagramData", "data1.xml"),
                    ("diagramLayout", "layout1.xml"),
                    ("diagramQuickStyle", "quickStyle1.xml"),
                    ("diagramColors", "colors1.xml"),
                ), start=1)
            )
            rels = rels.replace("</Relationships>", f"{extra_rels}</Relationships>")
            contents["ppt/slides/_rels/slide7.xml.rels"] = rels.encode("utf-8")

            slide7 = contents["ppt/slides/slide7.xml"].decode("utf-8")
            graphic_frame = (
                '<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="900" name="SmartArt 1"/>'
                '<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>'
                '<p:xfrm><a:off x="457200" y="1600200"/><a:ext cx="8229600" cy="4525963"/></p:xfrm>'
                '<a:graphic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                '<a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
                '<dgm:relIds xmlns:dgm="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
                'r:dm="rIdDg1" r:lo="rIdDg2" r:qs="rIdDg3" r:cs="rIdDg4"/>'
                '</a:graphicData></a:graphic></p:graphicFrame>')
            slide7 = slide7.replace("</p:spTree>", f"{graphic_frame}</p:spTree>", 1)
            contents["ppt/slides/slide7.xml"] = slide7.encode("utf-8")

            presentation_rels = contents["ppt/_rels/presentation.xml.rels"].decode("utf-8")
            rel_to_slide = {
                target: rel_id
                for rel_id, target in re.findall(
                    r'Id="([^"]+)"[^>]*Target="([^"]+)"', presentation_rels)
            }

            def rid(slide_no: int) -> str:
                return rel_to_slide[f"slides/slide{slide_no}.xml"]

            shows_xml = (
                '<p14:custShowLst xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
                '<p14:custShow name="放映方案1" id="0"><p14:sldLst>'
                + "".join(f'<p14:sldId r:id="{rid(n)}"/>' for n in (1, 2, 4, 7))
                + '</p14:sldLst></p14:custShow>'
                '<p14:custShow name="放映方案2" id="1"><p14:sldLst>'
                + "".join(f'<p14:sldId r:id="{rid(n)}"/>' for n in (1, 2, 3, 5, 6))
                + "</p14:sldLst></p14:custShow></p14:custShowLst>")
            presentation = contents["ppt/presentation.xml"].decode("utf-8")
            if "<p:extLst>" in presentation:
                presentation = presentation.replace("<p:extLst>", f"{shows_xml}<p:extLst>", 1)
            else:
                presentation = presentation.replace(
                    "</p:presentation>", f"{shows_xml}</p:presentation>", 1)
            contents["ppt/presentation.xml"] = presentation.encode("utf-8")

            content_types = contents["[Content_Types].xml"].decode("utf-8")
            overrides = "".join(
                f'<Override PartName="/ppt/diagrams/{part}" ContentType="{ctype}"/>'
                for part, ctype in (
                    ("data1.xml",
                     "application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml"),
                    ("layout1.xml",
                     "application/vnd.openxmlformats-officedocument.drawingml.diagramLayout+xml"),
                    ("quickStyle1.xml",
                     "application/vnd.openxmlformats-officedocument.drawingml.diagramStyle+xml"),
                    ("colors1.xml",
                     "application/vnd.openxmlformats-officedocument.drawingml.diagramColors+xml"),
                )
            )
            content_types = content_types.replace("</Types>", f"{overrides}</Types>")
            contents["[Content_Types].xml"] = content_types.encode("utf-8")

        _zip_patch(path, mutate)

    patch_pptx(pptx_path)


# ---------------------------------------------------------------- paper02
PNG_1PX = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx"
    b"\x9cc\xf8\xcf\xc0\xf0\x1f\x00\x05\x05\x02\x00_\xc8\xf2\xf5\x00\x00"
    b"\x00\x00IEND\xaeB`\x82"
)


def make_paper02(out: Path) -> None:
    from agentbench.ncre_assets.exam_data_paper02 import (
        GRADE_ROWS,
        PPT_KEYWORDS,
        PPT_ORG,
        PPT_OUTLINE,
        PPT_TITLE,
        SCHEDULE_HEADER,
        SCHEDULE_ROWS,
        SIGNUP_STEPS,
        SPEAKER_INTRO,
        SUBJECTS,
        WORD_PAGE2_TITLE,
        WORD_POSTER_PAGE1,
    )
    from openpyxl.chart import BarChart, Reference

    answers = {
        "q01": "A", "q02": "A", "q03": "A", "q04": "A", "q05": "A",
        "q06": "A", "q07": "C", "q08": "B", "q09": "D", "q10": "B",
        "q11": "A", "q12": "B", "q13": "C", "q14": "D", "q15": "B",
        "q16": "A", "q17": "A", "q18": "A", "q19": "C", "q20": "A",
    }
    (out / "answers.json").write_text(json.dumps(answers), encoding="utf-8")

    # ---- WORD.docx 海报
    doc = Document()
    section = doc.sections[0]
    section.page_width = Cm(27)
    section.page_height = Cm(35)
    section.top_margin = Cm(5)
    section.bottom_margin = Cm(5)
    section.left_margin = Cm(3)
    section.right_margin = Cm(3)

    title_paragraph = doc.add_paragraph()
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_paragraph.add_run("领慧讲堂")
    title_run.font.name = "微软雅黑"
    title_run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    title_run.font.size = Pt(62)
    title_run.font.color.rgb = RGBColor(0xFF, 0, 0)
    for text in WORD_POSTER_PAGE1[1:]:
        doc.add_paragraph(text.replace("报告人：", "报告人：赵蕈"))
    for index, step in enumerate(SIGNUP_STEPS, start=1):
        doc.add_paragraph(f"报名流程{index}：{step}")
    intro = doc.add_paragraph(SPEAKER_INTRO)
    intro.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    intro.paragraph_format.first_line_indent = Pt(21)

    page2 = doc.add_section(WD_SECTION_START.NEW_PAGE)
    page2.orientation = WD_ORIENT.LANDSCAPE
    page2.page_width = Cm(29.7)
    page2.page_height = Cm(21)
    doc.add_paragraph(WORD_PAGE2_TITLE)
    schedule = doc.add_table(rows=5, cols=3)
    for column, header in enumerate(SCHEDULE_HEADER):
        schedule.rows[0].cells[column].text = header
    for row_index, row_values in enumerate(SCHEDULE_ROWS, start=1):
        for column, value in enumerate(row_values):
            schedule.rows[row_index].cells[column].text = str(value)
    word_path = out / "WORD.docx"
    doc.save(word_path)
    _add_background(word_path, "FFF2CC")

    # ---- Excel.xlsx 成绩单
    workbook = openpyxl.load_workbook(io.BytesIO(build_assets.build_paper02_excel_asset()))
    grades = workbook["成绩表"]
    red_fill = PatternFill(start_color="FFFFC7CE", end_color="FFFFC7CE", fill_type="solid")
    blue_font = Font(color="FF0000FF")
    for row in range(2, 14):
        grades.cell(row=row, column=3).value = f'=MID(A{row},3,2)&"班"'
        grades.cell(row=row, column=11).value = f"=SUM(D{row}:J{row})"
        avg_cell = grades.cell(row=row, column=12)
        avg_cell.value = f"=AVERAGE(D{row}:J{row})"
        avg_cell.number_format = "0.00"
        for column in range(4, 11):
            grades.cell(row=row, column=column).number_format = "0.00"
    grades.conditional_formatting.add(
        "D2:F13", CellIsRule(operator="greaterThanOrEqual", formula=["110"], fill=red_fill))
    grades.conditional_formatting.add(
        "G2:J13", CellIsRule(operator="greaterThan", formula=["95"], font=blue_font))

    summary = workbook.create_sheet("分类汇总")
    summary.sheet_properties.tabColor = "4472C4"
    # 分类汇总：按班级汇总（e6 检查班级均值）+ 转置分析数据（柱状图按班级分系列）
    by_class: dict[int, list] = {}
    for _student_id, _name, class_no, *scores in GRADE_ROWS:
        by_class.setdefault(class_no, []).append(scores)
    subject_averages = [
        [round(sum(group[subject_index] for group in groups) / len(groups), 2)
         for subject_index in range(len(SUBJECTS))]
        for groups in (by_class[1], by_class[2], by_class[3])
    ]
    class_names = ["1班", "2班", "3班"]
    summary["A1"] = "科目"
    for column, class_name in enumerate(class_names, start=2):
        summary.cell(row=1, column=column, value=class_name)
    for offset, subject in enumerate(SUBJECTS):
        row = 2 + offset
        summary.cell(row=row, column=1, value=subject)
        for column, averages in enumerate(subject_averages, start=2):
            summary.cell(row=row, column=column, value=averages[offset])
    summary.cell(row=9, column=1, value="班级均值对照")
    for offset, (class_name, averages) in enumerate(
            zip(class_names, subject_averages), start=10):
        summary.cell(row=offset, column=1, value=class_name)
        for column, value in enumerate(averages, start=2):
            summary.cell(row=offset, column=column, value=value)

    chart_sheet = workbook.create_sheet("柱状分析图")
    chart = BarChart()
    chart.type = "col"
    chart.grouping = "clustered"
    chart.title = "各班7科平均分（分析图）"
    data = Reference(summary, min_col=2, max_col=4, min_row=1, max_row=8)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(Reference(summary, min_col=1, min_row=2, max_row=8))
    chart_sheet.add_chart(chart, "A1")
    workbook.save(out / "Excel.xlsx")

    # ---- 水资源利用与节水.pptx
    prs = Presentation()
    layouts = prs.slide_layouts
    slide1 = prs.slides.add_slide(layouts[0])
    slide1.shapes.title.text_frame.text = PPT_TITLE
    box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1.5))
    frame = box.text_frame
    frame.text = PPT_ORG
    frame.add_paragraph().text = "2013年4月29日"

    heading_map = {"Heading 1": layouts[2], "Heading 2": layouts[1], "Normal": layouts[1]}
    current = None
    for style, text in PPT_OUTLINE:
        if style == "Heading 1" or current is None:
            current = prs.slides.add_slide(heading_map[style])
            current.shapes.title.text_frame.text = text
            continue
        body = next(p for p in current.placeholders if p.placeholder_format.idx != 0)
        if style == "Heading 2":
            paragraph = body.text_frame.add_paragraph()
            paragraph.level = 0
            paragraph.text = text
        else:
            paragraph = body.text_frame.add_paragraph()
            paragraph.level = 1
            paragraph.text = text

    end_slide = prs.slides.add_slide(layouts[1])
    end_slide.shapes.title.text_frame.text = "谢谢观看"

    slides = list(prs.slides)
    for target_slide in (slides[1], slides[3]):
        picture = target_slide.shapes.add_picture(
            io.BytesIO(PNG_1PX), Inches(8.0), Inches(6.0), Inches(0.8), Inches(0.8))
        picture.name = f"pic-{target_slide.slide_id}"
    link_box = slides[2].shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.5))
    link_run = link_box.text_frame.paragraphs[0].add_run()
    link_run.text = "更多节水知识：" + "、".join(PPT_KEYWORDS)
    link_run.hyperlink.address = "http://www.bjwater.gov.cn/"
    link_box2 = slides[3].shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.5))
    link_run2 = link_box2.text_frame.paragraphs[0].add_run()
    link_run2.text = "全国节水办公室"
    link_run2.hyperlink.address = "http://www.mwr.gov.cn/"

    pptx_path = out / "水资源利用与节水.pptx"
    prs.save(pptx_path)
    _zip_patch(pptx_path, lambda contents: contents.update({
        "ppt/theme/theme1.xml": re.sub(
            r'(<a:theme[^>]*name=")[^"]*(")', r"\1节水主题\2",
            contents["ppt/theme/theme1.xml"].decode("utf-8")).encode("utf-8"),
    }))


# ---------------------------------------------------------------- paper03
def make_paper03(out: Path) -> None:
    from agentbench.ncre_assets.exam_data_paper03 import (
        CENSUS_ROWS,
        COMPARE_TABLE,
        CONSULT_HEADER,
        CONSULT_ROWS,
        COURSEWARE_A_SLIDES,
        COURSEWARE_B_SLIDES,
        ENDING_TEXT,
        EXPECTED_GROWTH_DESC_GT50M,
        FOOTNOTE_TEXT,
        LINK_TARGET,
        LINK_TEXT,
        NEW_SLIDE_COMPARE_TITLE,
        NEW_SLIDE_TITLE_ONLY,
        PPT_FOOTER_TEXT,
        REPORT_BODY,
        REPORT_COVER,
        REPORT_INTRO,
    )
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls
    from lxml import etree
    from pptx.util import Pt as PptxPt

    answers = {
        "q01": "B", "q02": "B", "q03": "C", "q04": "C", "q05": "C",
        "q06": "B", "q07": "C", "q08": "B", "q09": "D", "q10": "A",
        "q11": "C", "q12": "B", "q13": "B", "q14": "B", "q15": "D",
        "q16": "A", "q17": "B", "q18": "A", "q19": "B", "q20": "A",
    }
    (out / "answers.json").write_text(json.dumps(answers), encoding="utf-8")

    # ---- 年报.docx
    doc = Document()
    section = doc.sections[0]
    section.page_width = Cm(18.4)
    section.page_height = Cm(26)
    section.top_margin = Cm(3.2)
    section.bottom_margin = Cm(3)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    doc.add_paragraph(REPORT_COVER[0].replace(" ", ""))
    for text in REPORT_COVER[1:]:
        doc.add_paragraph(text.replace(" ", ""))
    doc.paragraphs[2].add_run().add_break(WD_BREAK.PAGE)

    toc_paragraph = doc.add_paragraph()
    toc_paragraph._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="begin"/></w:r>'))
    toc_paragraph._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:instrText xml:space="preserve"> TOC \\o "1-3" \\h \\z \\u '
        '</w:instrText></w:r>'))
    toc_paragraph._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="end"/></w:r>'))

    doc.add_paragraph(REPORT_INTRO[0].replace(" ", ""))
    intro2_text = REPORT_INTRO[1].replace(" ", "")
    before_link, after_link = intro2_text.split(LINK_TEXT)
    intro2 = doc.add_paragraph(before_link)
    intro2.add_run(LINK_TEXT)
    intro2.add_run(after_link)

    style_map = {"h1": "Heading 1", "h2": "Heading 2", "h3": "Heading 3"}
    blue_paragraph = None
    for kind, text in REPORT_BODY:
        text = text.replace(" ", "")
        if kind == "blue":
            blue_paragraph = doc.add_paragraph(text)
        else:
            paragraph = doc.add_paragraph(text)
            if kind in style_map:
                paragraph.style = doc.styles[style_map[kind]]

    # 咨询情况表（替换蓝色段）+ 饼图
    table = doc.add_table(rows=5, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for column, header in enumerate(CONSULT_HEADER):
        table.rows[0].cells[column].text = header
    for row_index, (mode, count, ratio) in enumerate(CONSULT_ROWS, start=1):
        table.rows[row_index].cells[0].text = mode
        table.rows[row_index].cells[1].text = str(count)
        table.rows[row_index].cells[2].text = f"{ratio:g}"
    blue_paragraph._p.addnext(table._tbl)
    blue_paragraph._p.getparent().remove(blue_paragraph._p)

    # 咨询情况饼图：openpyxl 生成 chart 部件后手工注入 docx 包
    # （当前 python-docx 版本无 docx.chart，无法 Document.add_chart）
    def add_pie_chart(docx_path: Path, anchor_paragraph) -> None:
        chart_wb = openpyxl.Workbook()
        chart_ws = chart_wb.active
        chart_ws["A1"] = "咨询方式"
        chart_ws["B1"] = "人次"
        for offset, (mode, count, _ratio) in enumerate(CONSULT_ROWS[:-1], start=2):
            chart_ws.cell(row=offset, column=1, value=mode)
            chart_ws.cell(row=offset, column=2, value=count)
        from openpyxl.chart import PieChart, Reference
        from openpyxl.chart.label import DataLabelList

        pie = PieChart()
        pie.add_data(Reference(chart_ws, min_col=2, min_row=1, max_row=4),
                     titles_from_data=True)
        pie.set_categories(Reference(chart_ws, min_col=1, min_row=2, max_row=4))
        pie.dataLabels = DataLabelList(showPercent=True, showVal=False)
        chart_ws.add_chart(pie, "D2")
        chart_buffer = io.BytesIO()
        chart_wb.save(chart_buffer)
        with zipfile.ZipFile(io.BytesIO(chart_buffer.getvalue())) as archive:
            chart_xml = archive.read("xl/charts/chart1.xml")

        def mutate(contents):
            contents["word/charts/chart1.xml"] = chart_xml
            drawing_xml = (
                '<w:drawing xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
                'xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" '
                'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                'xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                '<wp:inline distT="0" distB="0" distL="0" distR="0">'
                '<wp:extent cx="4572000" cy="2743200"/><wp:effectExtent l="0" t="0" r="0" b="0"/>'
                '<wp:docPr id="100" name="咨询情况饼图"/>'
                '<a:graphic><a:graphicData '
                'uri="http://schemas.openxmlformats.org/drawingml/2006/chart">'
                '<c:chart r:id="rIdChart1"/></a:graphicData></a:graphic>'
                '</wp:inline></w:drawing>')
            anchor_paragraph._p.append(parse_xml(
                f'<w:r {nsdecls("w")}>{drawing_xml}</w:r>'))
            contents["word/document.xml"] = etree.tostring(
                anchor_paragraph._p.getroottree().getroot(),
                xml_declaration=True, encoding="UTF-8", standalone=True)
            chart_rels = 'word/_rels/document.xml.rels'
            rels_xml = contents[chart_rels].decode("utf-8")
            rels_xml = rels_xml.replace(
                "</Relationships>",
                '<Relationship Id="rIdChart1" Type="http://schemas.openxmlformats.org/'
                'officeDocument/2006/relationships/chart" '
                'Target="charts/chart1.xml"/></Relationships>')
            contents[chart_rels] = rels_xml.encode("utf-8")
            content_types = contents["[Content_Types].xml"].decode("utf-8")
            if "word/charts/chart1.xml" not in content_types:
                content_types = content_types.replace(
                    "</Types>",
                    '<Override PartName="/word/charts/chart1.xml" ContentType="application/'
                    'vnd.openxmlformats-officedocument.drawingml.chart+xml"/></Types>')
                contents["[Content_Types].xml"] = content_types.encode("utf-8")

        _zip_patch(docx_path, mutate)

    # 奇偶页眉与页码
    section.different_odd_even = True
    odd_header = section.header
    odd_header.paragraphs[0].text = "北京市政府信息公开工作年度报告"
    odd_header.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    odd_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="begin"/></w:r>'))
    odd_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:instrText xml:space="preserve"> PAGE </w:instrText></w:r>'))
    odd_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="end"/></w:r>'))
    even_header = section.even_page_header
    even_header.paragraphs[0].text = "北京市政府信息公开工作年度报告"
    even_header.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT
    even_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="begin"/></w:r>'))
    even_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:instrText xml:space="preserve"> PAGE </w:instrText></w:r>'))
    even_header.paragraphs[0]._p.append(parse_xml(
        f'<w:r {nsdecls("w")}><w:fldChar w:fldCharType="end"/></w:r>'))

    intro2.runs[-1]._r.append(parse_xml(
        f'<w:footnoteReference w:id="1" {nsdecls("w")}/>'))

    report_path = out / "年报.docx"
    doc.save(report_path)
    add_pie_chart(report_path, doc.paragraphs[-1])

    def patch_report(contents):
        document_xml = contents["word/document.xml"].decode("utf-8")
        link_run_xml = (
            '<w:hyperlink r:id="rIdExtLink1"><w:r><w:rPr><w:color w:val="FF0000"/></w:rPr>'
            f'<w:t>{LINK_TEXT}</w:t></w:r></w:hyperlink>')
        document_xml = re.sub(
            r'<w:r>(?:(?!</w:r>).)*?<w:t[^>]*>' + re.escape(LINK_TEXT)
            + r'</w:t></w:r>', link_run_xml, document_xml, count=1)
        document_xml = document_xml.replace(
            "</w:sectPr>", '<w:cols w:num="2"/></w:sectPr>', 1)
        contents["word/document.xml"] = document_xml.encode("utf-8")

        rels_xml = contents["word/_rels/document.xml.rels"].decode("utf-8")
        rels_xml = rels_xml.replace(
            "</Relationships>",
            f'<Relationship Id="rIdExtLink1" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" '
            f'Target="{LINK_TARGET}" TargetMode="External"/></Relationships>')
        contents["word/_rels/document.xml.rels"] = rels_xml.encode("utf-8")

        footnote_reference = (
            '<w:footnote w:id="1"><w:p><w:r><w:footnoteRef/></w:r>'
            f'<w:r><w:t xml:space="preserve">{FOOTNOTE_TEXT}</w:t></w:r></w:p></w:footnote>')
        contents["word/footnotes.xml"] = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<w:footnote w:type="separator" w:id="-1"><w:p><w:r><w:separator/></w:r></w:p></w:footnote>'
            '<w:footnote w:type="continuationSeparator" w:id="0"><w:p><w:r>'
            '<w:continuationSeparator/></w:r></w:p></w:footnote>'
            + footnote_reference + '</w:footnotes>').encode("utf-8")

        settings_xml = contents["word/settings.xml"].decode("utf-8")
        if "<w:evenAndOddHeaders" not in settings_xml:
            settings_xml = re.sub(
                r"(<w:settings[^>]*>)", r"\1<w:evenAndOddHeaders/>", settings_xml, count=1)
        settings_xml = settings_xml.replace(
            "<w:evenAndOddHeaders/>", "<w:evenAndOddHeaders/><w:footnotePr/>", 1)
        contents["word/settings.xml"] = settings_xml.encode("utf-8")

        content_types = contents["[Content_Types].xml"].decode("utf-8")
        if "word/footnotes.xml" not in content_types:
            content_types = content_types.replace(
                "</Types>",
                '<Override PartName="/word/footnotes.xml" ContentType="application/'
                'vnd.openxmlformats-officedocument.wordprocessingml.footnotes+xml"/></Types>')
            contents["[Content_Types].xml"] = content_types.encode("utf-8")

    _zip_patch(report_path, patch_report)

    # ---- Excel.xlsx 人口普查
    workbook = openpyxl.Workbook()
    s5 = workbook.active
    s5.title = "第五次"
    s6 = workbook.create_sheet("第六次")
    census = {name: (pop2000, pop2010) for name, pop2000, pop2010 in CENSUS_ROWS}
    total_31_2000 = sum(p5 for p5, _ in census.values())
    total_31_2010 = sum(p6 for _, p6 in census.values())
    total_official_2000 = 1295330000
    total_official_2010 = 1339724852
    for sheet, year_index, total in ((s5, 0, total_official_2000),
                                     (s6, 1, total_official_2010)):
        sheet["A1"] = "地区"
        sheet["B1"] = "人口数(人)"
        sheet["C1"] = "比重(%)"
        for offset, (name, pop2000, pop2010) in enumerate(CENSUS_ROWS):
            row = 2 + offset
            population = pop2000 if year_index == 0 else pop2010
            sheet.cell(row=row, column=1, value=name)
            sheet.cell(row=row, column=2, value=population).number_format = "#,##0"
            sheet.cell(row=row, column=3,
                       value=round(population / total * 100, 2)).number_format = "0.00"
        sheet.add_table(Table(displayName=f"Tbl{sheet.title}", ref="A1:C32"))
        sheet.tables[f"Tbl{sheet.title}"].tableStyleInfo = TableStyleInfo(
            name="TableStyleMedium9", showRowStripes=True)

    compare = workbook.create_sheet("比较数据")
    for column, header in enumerate(
            ["地区", "2000年人口(人)", "2010年人口(人)", "增长数(人)", "比重变化(百分点)"], start=1):
        compare.cell(row=1, column=column, value=header)
    for offset, name in enumerate(sorted(census)):
        pop2000, pop2010 = census[name]
        row = 2 + offset
        compare.cell(row=row, column=1, value=name)
        compare.cell(row=row, column=2, value=pop2000).number_format = "#,##0"
        compare.cell(row=row, column=3, value=pop2010).number_format = "#,##0"
        compare.cell(row=row, column=4, value=pop2010 - pop2000).number_format = "#,##0"
        change = (pop2010 / total_official_2010 - pop2000 / total_official_2000) * 100
        compare.cell(row=row, column=5, value=round(change, 2)).number_format = "0.00"

    stats = workbook.create_sheet("统计指标")
    big_populations = [item[1] for item in EXPECTED_GROWTH_DESC_GT50M]
    for offset, (label, value) in enumerate([
        ("地区数", 31),
        ("2000年全国合计(31地区)", total_31_2000),
        ("2010年全国合计(31地区)", total_31_2010),
        ("超5000万地区数", len(big_populations)),
        ("超5000万地区2010年人口平均", round(sum(big_populations) / len(big_populations), 1)),
    ]):
        stats.cell(row=1 + offset, column=1, value=label)
        cell = stats.cell(row=1 + offset, column=2, value=value)
        cell.number_format = "#,##0.0" if isinstance(value, float) else "#,##0"

    pivot = workbook.create_sheet("透视汇总")
    for column, header in enumerate(["地区", "2010年人口(人)", "增长数(人)"], start=1):
        pivot.cell(row=1, column=column, value=header)
    for offset, (name, population, growth) in enumerate(EXPECTED_GROWTH_DESC_GT50M):
        row = 2 + offset
        pivot.cell(row=row, column=1, value=name)
        pivot.cell(row=row, column=2, value=population).number_format = "#,##0"
        pivot.cell(row=row, column=3, value=growth).number_format = "#,##0"
    pivot.cell(row=12, column=1, value="全国合计(31地区)")
    pivot.cell(row=12, column=2, value=total_31_2000).number_format = "#,##0"
    pivot.cell(row=12, column=3, value=total_31_2010).number_format = "#,##0"
    workbook.save(out / "Excel.xlsx")

    # ---- 物理课件.pptx（双主题合并）
    prs = Presentation()
    layouts = prs.slide_layouts
    title_layout, content_layout, title_only_layout = layouts[0], layouts[1], layouts[5]

    a_slides = []
    for kind, title, bullets in COURSEWARE_A_SLIDES:
        slide = prs.slides.add_slide(title_layout if kind == "title" else content_layout)
        slide.shapes.title.text_frame.text = title
        if kind != "title":
            body = next(p for p in slide.placeholders if p.placeholder_format.idx != 0)
            body.text_frame.clear()
            for offset, line in enumerate(bullets):
                paragraph = body.text_frame.paragraphs[0] if offset == 0 \
                    else body.text_frame.add_paragraph()
                paragraph.text = line
        a_slides.append(slide)

    new_slide = prs.slides.add_slide(title_only_layout)
    new_slide.shapes.title.text_frame.text = NEW_SLIDE_TITLE_ONLY

    def add_b_slide(kind, title, bullets):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text_frame.text = title
        body = next(p for p in slide.placeholders if p.placeholder_format.idx != 0)
        body.text_frame.clear()
        for offset, line in enumerate(bullets):
            paragraph = body.text_frame.paragraphs[0] if offset == 0 \
                else body.text_frame.add_paragraph()
            paragraph.text = line
        return slide

    for kind, title, bullets in COURSEWARE_B_SLIDES[:2]:
        add_b_slide(kind, title, bullets)

    compare_slide = prs.slides.add_slide(content_layout)
    compare_slide.shapes.title.text_frame.text = NEW_SLIDE_COMPARE_TITLE
    table_shape = compare_slide.shapes.add_table(
        4, 3, Inches(0.8), Inches(2.0), Inches(8.4), Inches(3.0))
    for row_index, row_values in enumerate(COMPARE_TABLE):
        for column, value in enumerate(row_values):
            table_shape.table.cell(row_index, column).text = value

    for kind, title, bullets in COURSEWARE_B_SLIDES[2:]:
        add_b_slide(kind, title, bullets)

    ending = prs.slides.add_slide(content_layout)
    ending.shapes.title.text_frame.text = ENDING_TEXT

    slides = list(prs.slides)  # [A0,A1,A2,新增,第3节,第4节,对比表,第5节,结尾]
    for text, source_index, target_index in (
        ("返回第2节熔化和凝固", 3, 2), ("返回第4节升华和凝华", 6, 5)):
        box = slides[source_index].shapes.add_textbox(
            Inches(0.5), Inches(6.8), Inches(6), Inches(0.5))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = PptxPt(14)
        box.click_action.target_slide = slides[target_index]

    pptx_path = out / "物理课件.pptx"
    prs.save(pptx_path)

    def patch_courseware(contents):
        slide_parts = sorted(
            (name for name in contents
             if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"\d+", name.rsplit("/", 1)[-1]).group()))
        # 页脚与幻灯片编号
        for index, part in enumerate(slide_parts):
            xml = contents[part].decode("utf-8")
            inject = (
                f'<p:sp><p:nvSpPr><p:cNvPr id="{900 + index}" name="页脚"/>'
                f'<p:cNvSpPr/><p:nvPr><p:ph type="ftr" sz="quarter"/></p:nvPr></p:nvSpPr>'
                f'<p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>'
                f'<a:p><a:r><a:t>{PPT_FOOTER_TEXT}</a:t></a:r></a:p></p:txBody></p:sp>')
            if index > 0:
                inject += (
                    f'<p:sp><p:nvSpPr><p:cNvPr id="{920 + index}" name="编号"/>'
                    f'<p:cNvSpPr/><p:nvPr><p:ph type="sldNum" sz="quarter"/></p:nvPr>'
                    f'</p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>'
                    f'<a:p/></p:txBody></p:sp>')
            xml = xml.replace("</p:spTree>", f"{inject}</p:spTree>", 1)
            xml = xml.replace("</p:cSld>",
                              '</p:cSld><p:transition spd="med"><p:fade/></p:transition>', 1)
            contents[part] = xml.encode("utf-8")
        # 第二主题（master2 + theme2 + 专用版式）
        master1 = contents["ppt/slideMasters/slideMaster1.xml"].decode("utf-8")
        master2 = re.sub(r'(<p:sldMasterIdLst>).*?(</p:sldMasterIdLst>)',
                         r'\1<p:sldMasterId id="2147483700" r:id="rId2"/>\2',
                         master1, flags=re.DOTALL)
        contents["ppt/slideMasters/slideMaster2.xml"] = master2.encode("utf-8")
        contents["ppt/slideMasters/_rels/slideMaster2.xml.rels"] = (
            b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/'
            b'2006/relationships/theme" Target="../theme/theme2.xml"/>'
            b'<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/'
            b'2006/relationships/slideLayout" Target="../slideLayouts/slideLayout50.xml"/>'
            b'</Relationships>')
        contents["ppt/theme/theme2.xml"] = re.sub(
            r'(<a:theme[^>]*name=")[^"]*(")', r"\1物理课件主题\2",
            contents["ppt/theme/theme1.xml"].decode("utf-8")).encode("utf-8")
        layout_xml = contents["ppt/slideLayouts/slideLayout6.xml"].decode("utf-8")
        contents["ppt/slideLayouts/slideLayout50.xml"] = layout_xml.encode("utf-8")
        contents["ppt/slideLayouts/_rels/slideLayout50.xml.rels"] = (
            b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/'
            b'2006/relationships/slideMaster" Target="../slideMasters/slideMaster2.xml"/>'
            b'</Relationships>')
        # 第4-9张改用第二主题版式
        for part in slide_parts[3:]:
            rels_name = part.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
            rels = contents[rels_name].decode("utf-8")
            rels = re.sub(r'(Type="[^"]*relationships/slideLayout"[^>]*Target=")[^"]*(")',
                          r"\g<1>../slideLayouts/slideLayout50.xml\g<2>", rels)
            contents[rels_name] = rels.encode("utf-8")
        # presentation.xml 与 rels、Content_Types 登记
        presentation = contents["ppt/presentation.xml"].decode("utf-8")
        presentation = presentation.replace(
            "</p:sldMasterIdLst>",
            '<p:sldMasterId id="2147483649" r:id="rIdMaster2"/></p:sldMasterIdLst>', 1)
        contents["ppt/presentation.xml"] = presentation.encode("utf-8")
        presentation_rels = contents["ppt/_rels/presentation.xml.rels"].decode("utf-8")
        presentation_rels = presentation_rels.replace(
            "</Relationships>",
            '<Relationship Id="rIdMaster2" Type="http://schemas.openxmlformats.org/'
            'officeDocument/2006/relationships/slideMaster" '
            'Target="slideMasters/slideMaster2.xml"/></Relationships>')
        contents["ppt/_rels/presentation.xml.rels"] = presentation_rels.encode("utf-8")
        content_types = contents["[Content_Types].xml"].decode("utf-8")
        content_types = content_types.replace(
            "</Types>",
            '<Override PartName="/ppt/slideMasters/slideMaster2.xml" ContentType="application/'
            'vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>'
            '<Override PartName="/ppt/slideLayouts/slideLayout50.xml" ContentType="application/'
            'vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
            '<Override PartName="/ppt/theme/theme2.xml" ContentType="application/'
            'vnd.openxmlformats-officedocument.theme+xml"/></Types>')
        contents["[Content_Types].xml"] = content_types.encode("utf-8")

    _zip_patch(pptx_path, patch_courseware)


# ---------------------------------------------------------------- 入口
BUILDERS = {
    "paper01": make_paper01,
    "paper02": make_paper02,
    "paper03": make_paper03,
}


def main(paper: str = "paper01") -> None:
    if paper not in BUILDERS:
        raise SystemExit(f"未知卷别: {paper}（可选 {'/'.join(BUILDERS)}）")
    out = DEV / ("reference-workspace" if paper == "paper01"
                 else f"reference-workspace-{paper}")
    if out.exists():
        shutil.rmtree(out)
    out.mkdir(parents=True)
    BUILDERS[paper](out)
    print("reference workspace ready:", out)


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "paper01")
